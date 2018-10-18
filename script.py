from pptx import Presentation
import csv

prs = Presentation()
title_slide_layout = prs.slide_layouts[8]

engwords = []
arwords = []
images = []
finalWords = []

with open('engwords.csv', 'rb') as csvfile:
    spamreader = csv.reader(csvfile, delimiter=',')
    for row in spamreader:
        for i in range(len(row)):
            engwords.append(row[i])
            # print engwords

with open('arwords.csv', 'rb') as csvfile:
    spamreader = csv.reader(csvfile, delimiter=',')
    for row in spamreader:
        for i in range(len(row)):
            arwords.append(row[i])
            # print arwords

with open('images.csv', 'rb') as csvfile:
    spamreader = csv.reader(csvfile, delimiter=',')
    for row in spamreader:
        for i in range(len(row)):
            images.append(row[i])
            # print images

for i in range(len(engwords)):
    finalWords.append(engwords[i] + "  -  " + arwords[i])


for i in range(len(images)):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    placeholder = slide.placeholders[1]
    print placeholder.placeholder_format.type
    title.text = finalWords[i]
    path = "images/" + images[i]
    picture = placeholder.insert_picture(path)


prs.save('final.pptx')
